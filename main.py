import re
import docx
from bs4 import BeautifulSoup
from pypdf import PdfReader

from sentence_transformers import SentenceTransformer, util

import warnings
import hdbscan

import numpy as np
import seaborn as sns

from transformers import BartTokenizer, BartForConditionalGeneration, BartConfig

import torch
from transformers import LongformerTokenizer, EncoderDecoderModel

import nltk
from nltk.corpus import stopwords
from nltk.tokenize import word_tokenize, sent_tokenize
from sklearn.feature_extraction.text import TfidfVectorizer
import matplotlib.pyplot as plt

nltk.download("punkt_tab")
nltk.download("stopwords")
nltk.download("punkt")

import matplotlib.pyplot as plt
import scipy.cluster.hierarchy as sch
from sklearn.metrics.pairwise import cosine_similarity

import plotly.express as px
from sklearn.manifold import TSNE

from wordcloud import WordCloud
import matplotlib.pyplot as plt

import pandas as pd
import json
import xml.etree.ElementTree as ET
import os
import warnings
import pptx

import io
from PIL import Image

warnings.filterwarnings("ignore")


def clean_text(text):
    text = re.sub(r"http\S+|www\S+|https\S+", "", text)
    text = re.sub(r"\s+", " ", text).strip()
    text = re.sub(r"[^\w\s,.]", "", text)
    return text


def extract_and_clean_text(file_path):
    text = ""
    if file_path.endswith(".docx"):
        doc = docx.Document(file_path)
        for paragraph in doc.paragraphs:
            text += paragraph.text + " "
    elif file_path.endswith(".txt"):
        with open(file_path, "r", encoding="utf-8") as f:
            text = f.read()
    elif file_path.endswith((".html", ".htm")):
        with open(file_path, "r", encoding="utf-8") as f:
            html_content = f.read()
        soup = BeautifulSoup(html_content, "html.parser")
        text = soup.get_text(separator=" ", strip=True)
    elif file_path.endswith(".pdf"):
        reader = PdfReader(file_path)
        for page in reader.pages:
            text += page.extract_text() + " "
    elif file_path.endswith(".csv"):
        df = pd.read_csv(file_path)
        text = " ".join(df.astype(str).agg(" ".join, axis=1))
    elif file_path.endswith(".xlsx"):
        df = pd.read_excel(file_path)
        text = " ".join(df.astype(str).agg(" ".join, axis=1))
    elif file_path.endswith(".json"):
        with open(file_path, "r", encoding="utf-8") as f:
            data = json.load(f)
        text = " ".join([str(item) for item in data])
    elif file_path.endswith(".xml"):
        tree = ET.parse(file_path)
        root = tree.getroot()
        text = " ".join([elem.text for elem in root.iter() if elem.text])
    elif file_path.endswith(".pptx"):
        from pptx import Presentation

        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + " "
    else:
        raise ValueError("Unsupported file type: {}".format(file_path))
    cleaned_text = clean_text(text)
    return cleaned_text


def clean_files(file_list):
    cleaned_files = []
    for file in file_list:
        cleaned_files.append(extract_and_clean_text(file))
    return cleaned_files


def get_embeddings(text):
    model = SentenceTransformer("all-mpnet-base-v2")
    embeddings = model.encode(text)
    return embeddings


def clustering_labels(embeddings):
    warnings.filterwarnings("ignore")
    embeddings = np.array(embeddings)
    if len(embeddings) < 2:
        raise ValueError(
            "Not enough data points for clustering. At least 2 are required."
        )
    min_cluster_size = min(2, len(embeddings))
    cluster = hdbscan.HDBSCAN(
        min_cluster_size=min_cluster_size,
        metric="euclidean",
        cluster_selection_method="eom",
    ).fit(embeddings)
    return cluster.labels_


def bart_summarizer(text):
    model_name_bart = "facebook/bart-large-cnn"
    tokenizer = BartTokenizer.from_pretrained(model_name_bart)
    model = BartForConditionalGeneration.from_pretrained(model_name_bart)
    tokenize_inputs = tokenizer.encode(
        text, return_tensors="pt", max_length=1024, truncation=True
    )
    ids_summarization = model.generate(
        tokenize_inputs, num_beams=4, max_length=150, early_stopping=True
    )
    summary_decoded = tokenizer.decode(ids_summarization[0], skip_special_tokens=True)
    return summary_decoded


def longformer_summarizer(text):
    tokenizer = LongformerTokenizer.from_pretrained("allenai/longformer-base-4096")
    model = EncoderDecoderModel.from_pretrained(
        "patrickvonplaten/longformer2roberta-cnn_dailymail-fp16"
    )
    inputs = tokenizer(
        text, return_tensors="pt", padding="longest", truncation=True
    ).input_ids
    ids_summarization = model.generate(inputs)
    summary_decoded = tokenizer.decode(ids_summarization[0], skip_special_tokens=True)
    return summary_decoded


def longformer_summarizer_long_text(
    text, max_chunk_length=4000, overlap=200, max_summary_length=1024
):
    tokenizer = LongformerTokenizer.from_pretrained("allenai/longformer-base-4096")
    model = EncoderDecoderModel.from_pretrained(
        "patrickvonplaten/longformer2roberta-cnn_dailymail-fp16"
    )
    device = torch.device("cuda" if torch.cuda.is_available() else "cpu")
    model = model.to(device)
    tokens = tokenizer.encode(text)
    if len(tokens) <= max_chunk_length:
        inputs = tokenizer(text, return_tensors="pt", padding="longest").input_ids.to(
            device
        )
        summary_ids = model.generate(inputs, max_length=max_summary_length)
        summary = tokenizer.decode(summary_ids[0], skip_special_tokens=True)
        return summary
    chunk_summaries = []
    for i in range(0, len(tokens), max_chunk_length - overlap):
        chunk_tokens = tokens[i : i + max_chunk_length]
        if len(chunk_tokens) < 100:
            continue
        chunk_text = tokenizer.decode(chunk_tokens, skip_special_tokens=True)
        inputs = tokenizer(
            chunk_text, return_tensors="pt", padding="longest"
        ).input_ids.to(device)
        summary_ids = model.generate(inputs, max_length=max_summary_length // 2)
        chunk_summary = tokenizer.decode(summary_ids[0], skip_special_tokens=True)
        chunk_summaries.append(chunk_summary)
    final_summary = " ".join(chunk_summaries)
    return final_summary


def summarize_text(text):
    bart_tokenizer = BartTokenizer.from_pretrained("facebook/bart-large-cnn")
    input_length = len(bart_tokenizer.encode(text))
    if input_length < 1024:
        summary = bart_summarizer(text)
    elif input_length < 4096:
        summary = longformer_summarizer(text)
    else:
        summary = longformer_summarizer_long_text(text)
    return summary


def summarize(embeddings, labels, cleaned_files):
    no_of_clusters = max(labels) + 1
    clusters_embeddings = []
    clusters_text = [""] * no_of_clusters
    for i in range(no_of_clusters):
        clusters_embeddings.append(embeddings[labels == i])
    noise_docs = []
    for label, text_chunk in zip(labels, cleaned_files):
        if label != -1:
            clusters_text[label] += text_chunk
        else:
            noise_docs.append(text_chunk)
    clusters_text.extend(noise_docs)
    cluster_texts_combined = ["".join(cluster) for cluster in clusters_text]
    final_summaries = [
        summarize_text(cluster_text) for cluster_text in cluster_texts_combined
    ]
    return final_summaries


def tfidf_plot(all_text):
    tokens = word_tokenize(all_text.lower())
    stop_words = set(stopwords.words("english"))
    filtered_tokens = [w for w in tokens if not w in stop_words and w.isalnum()]
    vectorizer = TfidfVectorizer()
    tfidf_matrix = vectorizer.fit_transform([" ".join(filtered_tokens)])
    feature_names = vectorizer.get_feature_names_out()
    tfidf_scores = tfidf_matrix.toarray()[0]
    top_n = 25
    top_indices = tfidf_scores.argsort()[-top_n:]
    top_words = [feature_names[i] for i in top_indices]
    top_scores = [tfidf_scores[i] for i in top_indices]
    fig, ax = plt.subplots(figsize=(10, 5))
    ax.barh(top_words, top_scores, color="skyblue")
    ax.set_xlabel("TF-IDF Score")
    ax.set_ylabel("Words")
    ax.set_title("Top {} Important Words (TF-IDF)".format(top_n))
    ax.invert_yaxis()
    return fig


def dendrogram_plot(embeddings, labels):
    similarity_matrix = cosine_similarity(embeddings)
    distance_matrix = 1 - similarity_matrix
    linkage_matrix = sch.linkage(distance_matrix, method="ward")
    dendrogram_labels = [
        f"Doc {i} (Cluster {labels[i]})" if labels[i] != -1 else f"Doc {i} (Noise)"
        for i in range(len(labels))
    ]
    fig, ax = plt.subplots(figsize=(12, 8))
    sch.dendrogram(
        linkage_matrix,
        labels=dendrogram_labels,
        orientation="right",
        leaf_font_size=10,
        ax=ax,
    )
    ax.set_title("Hierarchical Dendrogram of Document Clusters", fontsize=14)
    ax.set_xlabel("Distance", fontsize=12)
    ax.set_ylabel("Documents", fontsize=12)
    plt.tight_layout()
    return fig


def tsne_plot(embeddings, labels):
    n_samples = len(embeddings)
    if n_samples < 2:
        fig, ax = plt.subplots(figsize=(6, 4))
        ax.text(
            0.5,
            0.5,
            "t-SNE plot is not applicable for a single document.",
            fontsize=12,
            ha="center",
            va="center",
            wrap=True,
        )
        ax.axis("off")
        return fig
    perplexity = min(30, n_samples - 1)
    tsne = TSNE(n_components=2, perplexity=perplexity, random_state=42)
    reduced_embeddings = tsne.fit_transform(embeddings)
    fig, ax = plt.subplots(figsize=(8, 6))
    scatter = ax.scatter(
        reduced_embeddings[:, 0],
        reduced_embeddings[:, 1],
        c=labels,
        cmap="viridis",
        s=50,
        alpha=0.8,
    )
    ax.set_title("t-SNE Visualization of Document Clusters", fontsize=14)
    ax.set_xlabel("t-SNE Dimension 1", fontsize=12)
    ax.set_ylabel("t-SNE Dimension 2", fontsize=12)
    cbar = plt.colorbar(scatter, ax=ax)
    cbar.set_label("Cluster Labels", fontsize=12)
    return fig


def wordcloud_plot(all_text):
    wordcloud = WordCloud(width=800, height=400, background_color="white").generate(
        all_text
    )
    fig, ax = plt.subplots(figsize=(10, 5), facecolor=None)
    ax.imshow(wordcloud)
    ax.axis("off")
    plt.tight_layout(pad=0)
    buf = io.BytesIO()
    fig.savefig(buf, format="png")
    buf.seek(0)
    img = Image.open(buf)
    img_array = np.array(img)
    buf.close()
    plt.close(fig)
    return img_array


def summarize_docs(files_text):
    if files_text:
        cleaned_files = clean_files(files_text)
        if len(cleaned_files) == 1:
            summary = summarize_text(cleaned_files[0])
            return (
                f"Summary for the uploaded document:\n{summary}",
                None,
                None,
                None,
                None,
            )
        embeddings = get_embeddings(cleaned_files)
        if len(embeddings) < 2:
            return (
                "Not enough documents for clustering. Please upload more files.",
                None,
                None,
                None,
                None,
            )
        labels = clustering_labels(embeddings)
        summaries = summarize(embeddings, labels, cleaned_files)
        summary_output = "\n".join(
            [
                f"• Summary for cluster/doc {i+1}:\n{summary}"
                for i, summary in enumerate(summaries)
            ]
        )
        all_text = " ".join(cleaned_files)
        tfidf_fig = tfidf_plot(all_text)  # Get the tfidf plot figure
        dendrogram_fig = dendrogram_plot(
            embeddings, labels
        )  # Get the dendrogram plot figure
        tsne_fig = tsne_plot(embeddings, labels)  # Get the t-sne plot figure
        wordcloud_fig = wordcloud_plot(all_text)  # Get the wordcloud plot figure
        return summary_output, tfidf_fig, dendrogram_fig, tsne_fig, wordcloud_fig
    else:
        return "No files uploaded.", None, None, None, None


import gradio as gr

with gr.Blocks() as demo:
    gr.Markdown("# 📰 Multi-Document Summarization")

    with gr.Row():
        with gr.Column():
            file_upload = gr.Files(label="Upload Your Files")
            gr.Markdown(
                "### Supported File Types: 📄 `.docx` 📝 `.txt` 🌐 `.html` 📑 `.pdf` 📊 `.csv` 📈 `.xlsx` 🗂 `.json` 🗃 `.xml` 🎞 `.pptx`",
                elem_id="file-types-info",
            )
            summarize_btn = gr.Button("Summarize")

        with gr.Column():
            summary_output = gr.Textbox(label="• Bullet List of Summaries", lines=10)

    gr.Markdown("## 📊 Visualizations")

    with gr.Row():
        dendro = gr.Plot(label="Dendrogram")
        tsne = gr.Plot(label="t-SNE")

    with gr.Row():
        tfidf = gr.Plot(label="TF-IDF")

    with gr.Row():
        wordcloud = gr.Image(label="Word Cloud")

    summarize_btn.click(
        summarize_docs,
        inputs=file_upload,
        outputs=[summary_output, tfidf, dendro, tsne, wordcloud],
    )

demo.launch()
